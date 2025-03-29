
import streamlit as st
from pptx import Presentation
import os
from openai import OpenAI
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
from io import BytesIO
from fpdf import FPDF

st.set_page_config(page_title="Enhanced Persona From PPTX with Avatars", layout="wide")
st.title("🧠 Persona Generator from PowerPoint + DALL·E Avatars")

# Debug: Show environment check
st.markdown("🔍 Debug Mode: DALL·E + GPT Enhanced")

# Load OpenAI key from environment variable
client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])

# Upload PPTX file
uploaded_file = st.file_uploader("Upload a PowerPoint (.pptx) with segmentation analysis", type=["pptx"])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def generate_gpt_response(prompt):
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a senior market research strategist."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip()

def generate_dalle_image(description):
    dalle_response = client.images.generate(
        model="dall-e-3",
        prompt=description,
        size="1024x1024",
        quality="standard",
        n=1
    )
    return dalle_response.data[0].url

# Session storage
if "summary" not in st.session_state:
    st.session_state.summary = ""
if "personas" not in st.session_state:
    st.session_state.personas = ""
if "persona_blocks" not in st.session_state:
    st.session_state.persona_blocks = []

# Generate segmentation summary
if uploaded_file and st.button("🔍 Generate Segmentation Summary"):
    ppt_text = extract_text_from_pptx(uploaded_file)
    with st.expander("📄 Slide Text Extracted"):
        st.text(ppt_text[:2000])
    st.info("Sending content to GPT for strategic summary...")
    summary_prompt = (
        "You are SAMI AI, an advanced market insights engine. Analyze the following segmentation slides and "
        "produce a strategic summary. Include:
"
        "- Segment descriptions (demographics, attitudes)
"
        "- Key differentiators
"
        "- Strategic implications for acquisition, loyalty, and innovation

"
        f"Slides:
{ppt_text[:4000]}"
    )
    summary = generate_gpt_response(summary_prompt)
    st.session_state.summary = summary
    st.subheader("📌 Strategic Summary")
    st.markdown(summary)

# Generate personas
if st.session_state.summary and st.button("👥 Generate Personas"):
    persona_prompt = (
        "Based on this segmentation summary, generate 5 distinct personas. Each should include:
"
        "- Name
"
        "- Description
"
        "- Traits
"
        "- Pain Points
"
        "- Motivations
"
        "- Preferred Channels
"
        "- Example Messaging

"
        f"Summary:
{st.session_state.summary}"
    )
    personas = generate_gpt_response(persona_prompt)
    st.session_state.personas = personas
    st.subheader("🎯 Personas")
    st.markdown(personas)

    # Parse into name + block
    blocks = personas.split("Persona ")
    parsed = []
    for block in blocks:
        if ":" in block:
            name = block.split(":")[0].strip()
            desc = block.split("Description:")[1].split("\n")[0].strip() if "Description:" in block else ""
            parsed.append((name, desc))
    st.session_state.persona_blocks = parsed

# Generate avatars
if st.session_state.persona_blocks and st.button("🖼️ Generate Persona Avatars"):
    st.subheader("🖼️ Persona Avatars")
    for name, desc in st.session_state.persona_blocks:
        try:
            with st.spinner(f"Generating image for {name}..."):
                image_prompt = f"Professional digital portrait of {name}, representing traits of the persona. Style: realistic, clean background, business attire."
                url = generate_dalle_image(image_prompt)
                st.image(url, caption=name, use_container_width=True)
                st.success(f"✅ Avatar generated for {name}")
        except Exception as e:
            st.error(f"❌ Failed to generate image for {name}: {e}")

# Export to PDF
if st.session_state.personas and st.button("📄 Download PDF Summary"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=10)
    pdf.multi_cell(0, 5, "📌 Strategic Summary\n" + st.session_state.summary)
    pdf.ln(4)
    pdf.multi_cell(0, 5, "🎯 Personas\n" + st.session_state.personas)
    buffer = BytesIO()
    pdf.output(buffer, 'F')
    buffer.seek(0)
    st.download_button("📥 Download PDF", buffer, file_name="persona_report.pdf", mime="application/pdf")
