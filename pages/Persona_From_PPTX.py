import streamlit as st
from pptx import Presentation
import os
from openai import OpenAI
import re
from fpdf import FPDF
from io import BytesIO
import requests

st.set_page_config(page_title="🧠 Persona Generator", layout="wide")
st.title("🧠 Persona Generator from PowerPoint + DALL·E Avatars")

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

uploaded_file = st.file_uploader("Upload a PowerPoint (.pptx) with segmentation analysis", type=["pptx"])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def generate_gpt_response(prompt):
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a senior market research strategist."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip()

def generate_dalle_image(description):
    st.write(f"👀 Description: {description}")
st.write(f"🧠 Final Prompt to DALL·E: {safe_prompt}")
    dalle_response = client.images.generate(
        model="dall-e-3",
        prompt=description,
        size="1024x1024",
        quality="standard",
        n=1
    )
    return dalle_response.data[0].url

# Session State Init
if "summary" not in st.session_state:
    st.session_state.summary = ""
if "personas" not in st.session_state:
    st.session_state.personas = ""
if "avatar_urls" not in st.session_state:
    st.session_state.avatar_urls = {}
if "persona_blocks" not in st.session_state:
    st.session_state.persona_blocks = []

# Step 1: Upload + Generate Strategic Summary
if uploaded_file and st.button("🔍 Generate Segmentation Summary"):
    ppt_text = extract_text_from_pptx(uploaded_file)
    with st.expander("📄 Slide Text Extracted"):
        st.text(ppt_text[:2000])

    st.info("Sending content to GPT for strategic summary...")
    summary_prompt = f"""You are SAMI AI, an advanced market insights engine. Analyze the following segmentation slides and produce a strategic summary. Include:
- Segment descriptions (demographics, attitudes)
- Key differentiators
- Strategic implications for acquisition, loyalty, and innovation

Slides:
{ppt_text[:4000]}"""

    st.session_state.summary = generate_gpt_response(summary_prompt)
    st.subheader("📌 Strategic Summary")
    st.markdown(st.session_state.summary)

# Step 2: Generate Personas
if st.session_state.summary and st.button("👥 Generate Personas"):
    summary_text = st.session_state.summary
    match = re.search(r'(\d+)\s+segments?', summary_text.lower())
    num_segments = int(match.group(1)) if match else 5

    persona_prompt = f"""You are SAMI AI, an advanced segmentation strategist.

Based on the segmentation summary below, generate exactly {num_segments} distinct personas.

Each persona must include:
## Name
## Description
## Traits
## Pain Points
## Motivations
## Preferred Channels
## Example Messaging

Each persona should be clearly separated and fully written. Do not skip any. Make the names creative but realistic.

Segmentation Summary:
{summary_text}"""

    personas = generate_gpt_response(persona_prompt)
    st.session_state.personas = personas
    st.session_state.persona_blocks = personas.split("## Name")[1:]

    st.subheader("🎯 Personas")
    st.markdown(personas)

# Step 3: Generate Avatars (AFTER personas are shown)
# Step 3: Generate Avatars (AFTER personas are shown)
if st.session_state.persona_blocks and st.button("🎨 Generate Avatars"):
    st.session_state.avatar_urls = {}

    for block in st.session_state.persona_blocks:
        lines = block.strip().split("\n")
        name = lines[0].strip()
        description = ""
        if "## Description" in block:
            try:
                description = block.split("## Description")[1].split("\n")[0].strip()
            except:
                continue

        # Build safe, enhanced prompt
        if not description or len(description.strip()) < 15:
            safe_prompt = "Realistic business portrait of a professional adult with a clean background and confident expression"
        else:
            safe_prompt = f"Realistic professional portrait of a person described as: {description}"

        try:
            st.write(f"🧠 Prompt sent to DALL·E for {name}: `{safe_prompt}`")
            image_url = generate_dalle_image(safe_prompt)
            st.session_state.avatar_urls[name] = image_url
            st.success(f"✅ Avatar generated for: {name}")
        except Exception as e:
            st.error(f"❌ Error generating avatar for {name}: {e}")
            # Optional fallback placeholder
            st.session_state.avatar_urls[name] = "https://via.placeholder.com/256x256.png?text=No+Image"


# Step 4: Display Avatars with Download Buttons
if st.session_state.avatar_urls:
    st.subheader("🖼️ Persona Avatars")
    for name, url in st.session_state.avatar_urls.items():
        try:
            response = requests.get(url)
            if response.status_code == 200:
                image_bytes = BytesIO(response.content)
                st.image(image_bytes, caption=name, width=256)
                st.download_button(
                    label=f"⬇️ Download {name}'s Avatar",
                    data=image_bytes,
                    file_name=f"{name}_avatar.jpg",
                    mime="image/jpeg"
                )
        except Exception as e:
            st.warning(f"⚠️ Could not display/download avatar for {name}: {e}")

# Step 5: PDF Summary Export (No Images)
def clean_text(text):
    if isinstance(text, str):
        return text.encode("latin-1", "ignore").decode("latin-1")
    return text

if st.session_state.personas and st.button("📄 Download PDF Summary (Text Only)"):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=10)

    summary_text = clean_text("📌 Strategic Summary\n" + st.session_state.summary)
    pdf.multi_cell(0, 5, summary_text)
    pdf.ln(6)

    for block in st.session_state.persona_blocks:
        lines = block.strip().split("\n")
        name = lines[0].strip()
        pdf.set_font("Arial", style='B', size=11)
        pdf.cell(0, 10, clean_text(name), ln=True)
        pdf.set_font("Arial", size=10)
        persona_text = clean_text("\n".join(lines[1:]))
        pdf.multi_cell(0, 5, persona_text)
        pdf.ln(4)

    pdf_bytes = pdf.output(dest='S').encode('latin-1')
    pdf_buffer = BytesIO(pdf_bytes)

    st.download_button(
        label="📥 Download PDF Summary (No Images)",
        data=pdf_buffer,
        file_name="persona_report.pdf",
        mime="application/pdf"
    )
