import streamlit as st
from pptx import Presentation
import os
from openai import OpenAI
import re
from fpdf import FPDF
from io import BytesIO
import requests

st.set_page_config(page_title="🧠 Persona Generator", layout="wide")
st.title("🧠 Persona Generator from PowerPoint + DALL·E Avatars")

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

uploaded_file = st.file_uploader("Upload a PowerPoint (.pptx) with segmentation analysis", type=["pptx"])

# Session State Init
for key in ["summary", "personas", "avatar_urls", "persona_blocks"]:
    if key not in st.session_state:
        st.session_state[key] = "" if key in ["summary", "personas"] else {}

# --- Helper Functions ---
def extract_text_from_pptx(file):
    prs = Presentation(file)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]).strip()

def generate_gpt_response(prompt):
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a senior market research strategist."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip()

def generate_dalle_image(description):
    dalle_prompt = description.strip() if len(description.strip()) > 10 else "realistic portrait of a person based on persona description"
    response = client.images.generate(
        model="dall-e-3",
        prompt=dalle_prompt,
        size="1024x1024",
        quality="standard",
        n=1
    )
    return response.data[0].url

def clean_text(text):
    return text.encode("latin-1", "ignore").decode("latin-1") if isinstance(text, str) else ""

# --- Step 1: Upload & Generate Summary ---
if uploaded_file and st.button("🔍 Generate Segmentation Summary"):
    ppt_text = extract_text_from_pptx(uploaded_file)
    with st.expander("📄 Slide Text Extracted"):
        st.text(ppt_text[:2000])

    st.info("Sending content to GPT for strategic summary...")
    summary_prompt = f"""
    You are SAMI AI, an advanced market insights engine. Analyze the following segmentation slides and produce a strategic summary. Include:
    - Segment descriptions (demographics, attitudes)
    - Key differentiators
    - Strategic implications for acquisition, loyalty, and innovation

    Slides:
    {ppt_text[:4000]}
    """
    st.session_state.summary = generate_gpt_response(summary_prompt)
    st.subheader("📌 Strategic Summary")
    st.markdown(st.session_state.summary)

# --- Step 2: Generate Personas ---
if st.session_state.summary and st.button("👥 Generate Personas"):
    match = re.search(r'(\d+)\s+segments?', st.session_state.summary.lower())
    num_segments = int(match.group(1)) if match else 5

    persona_prompt = f"""You are SAMI AI, an advanced segmentation strategist.
Based on the summary below, generate exactly {num_segments} distinct personas.
Each persona must include:
## Name
## Description
## Traits
## Pain Points
## Motivations
## Preferred Channels
## Example Messaging

Segmentation Summary:
{st.session_state.summary}
"""
    result = generate_gpt_response(persona_prompt)
    st.session_state.personas = result
    st.session_state.persona_blocks = result.split("## Name")[1:]

    st.subheader("🎯 Personas")
    st.markdown(result)

# --- Step 3: Generate DALL·E Avatars ---
if st.session_state.persona_blocks and st.button("🎨 Generate Avatars"):
    for block in st.session_state.persona_blocks:
        lines = block.strip().split("\n")
        name = lines[0].strip()
        try:
            desc_match = re.search(r"## Description\n(.+?)(\n##|$)", block, re.DOTALL)
            description = desc_match.group(1).strip() if desc_match else "A professional individual in a corporate setting"
            st.write(f"🧪 Generating avatar for: {name}")
            image_url = generate_dalle_image(description)
            st.session_state.avatar_urls[name] = image_url
            st.success(f"✅ Avatar generated for: {name}")
        except Exception as e:
            st.error(f"❌ Error generating avatar for {name}: {e}")

# --- Step 4: Show Avatars
if st.session_state.avatar_urls:
    st.subheader("🖼️ Persona Avatars")
    for name, url in st.session_state.avatar_urls.items():
        try:
            response = requests.get(url)
            if response.status_code == 200:
                img_bytes = BytesIO(response.content)
                st.image(img_bytes, caption=name, width=256)
                st.download_button(f"⬇️ Download {name}'s Avatar", img_bytes, file_name=f"{name}_avatar.jpg", mime="image/jpeg")
        except Exception as e:
            st.warning(f"⚠️ Could not display/download avatar for {name}: {e}")

# --- Step 5: PDF Export
if st.session_state.personas and st.button("📄 Download PDF Summary (Text Only)"):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=10)

    pdf.multi_cell(0, 5, clean_text("📌 Strategic Summary\n" + st.session_state.summary))
    pdf.ln(6)

    for block in st.session_state.persona_blocks:
        lines = block.strip().split("\n")
        name = lines[0].strip()
        pdf.set_font("Arial", style='B', size=11)
        pdf.cell(0, 10, clean_text(name), ln=True)
        pdf.set_font("Arial", size=10)
        pdf.multi_cell(0, 5, clean_text("\n".join(lines[1:])))
        pdf.ln(4)

    pdf_bytes = pdf.output(dest='S').encode('latin-1')
    st.download_button("📥 Download PDF Summary (No Images)", BytesIO(pdf_bytes), file_name="persona_report.pdf", mime="application/pdf")
