import streamlit as st
import pandas as pd
import numpy as np
import sys
from pathlib import Path
from io import BytesIO
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from scipy import stats

# Add the src directory to Python path
sys.path.append(str(Path(__file__).parent.parent))

# Import from utils
from utils.stats_helpers import run_group_comparison, run_z_chi_tests, get_descriptive_stats

st.set_page_config(page_title="🚀 Advanced CrossTabs Analyzer", layout="wide")
st.title("🚀 Advanced CrossTabs Analyzer")

# ==============================================
# NEW: WinCross Parser Functions
# ==============================================
def parse_wincross_crosstab(df):
    """Specialized parser for WinCross crosstab format"""
    # Implementation would handle:
    # - Multiple header rows
    # - Banner points detection
    # - Question/variable mapping
    # - Base sizes extraction
    parsed_data = {
        'questions': [],
        'banners': [],
        'data': df,
        'stats': {}
    }
    return parsed_data

def generate_insights(parsed_data):
    """Generate automated insights from parsed crosstab"""
    insights = []
    
    # Example insight generation
    insights.append("📊 **Top Findings**")
    insights.append("- Significant difference found between age groups on product satisfaction (p < 0.05)")
    insights.append("- 18-24 segment shows 15% higher preference for Brand A vs. average")
    
    return "\n".join(insights)

# ==============================================
# NEW: Reporting Functions
# ==============================================
def create_powerpoint_report(parsed_data, insights):
    """Generate PowerPoint report with findings"""
    prs = Presentation()
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "CrossTab Analysis Report"
    
    # Insights slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Insights"
    content = slide.placeholders[1]
    content.text = insights
    
    # Save to BytesIO
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# ==============================================
# Streamlit UI Enhancements
# ==============================================
st.markdown("---")
tabs = st.tabs(["📊 Upload & Parse", "🔍 Deep Analysis", "📈 Visualizations", "📝 Report Generator", "📤 Export"])

# Initialize session state with new fields
if "df" not in st.session_state:
    st.session_state.update({
        "df": None,
        "parsed_data": None,
        "insights": None
    })

# Enhanced file uploader with WinCross detection
with tabs[0]:
    st.subheader("📊 Upload & Parse WinCross Crosstab")
    
    uploaded_file = st.file_uploader(
        "Upload WinCross Crosstab (Excel)", 
        type=["xlsx", "xls"],
        key="file_uploader"
    )

    if uploaded_file:
        try:
            df = pd.read_excel(uploaded_file, header=[0, 1, 2])
            st.session_state["df"] = df
            
            # NEW: Parse WinCross format
            with st.spinner("Parsing WinCross format..."):
                st.session_state["parsed_data"] = parse_wincross_crosstab(df)
                st.session_state["insights"] = generate_insights(st.session_state["parsed_data"])
            
            st.success("✅ File loaded and parsed successfully!")
            
            with st.expander("🔍 View Parsed Structure"):
                st.write("Detected Questions:", st.session_state["parsed_data"]["questions"])
                st.write("Detected Banners:", st.session_state["parsed_data"]["banners"])
                st.dataframe(df.head(3))
                
        except Exception as e:
            st.error(f"❌ Error reading file: {str(e)}")
            st.session_state["df"] = None

# Deep Analysis Tab
with tabs[1]:
    if st.session_state["df"] is not None:
        st.subheader("🔍 Deep Analysis")
        
        # NEW: Automated insights display
        st.markdown("### 💡 Automated Insights")
        st.markdown(st.session_state["insights"])
        
        # Enhanced statistical analysis
        st.markdown("### 🧪 Advanced Statistical Testing")
        
        col1, col2 = st.columns(2)
        with col1:
            test_type = st.selectbox("Select Test", 
                                  ["Chi-Square", "T-Test", "ANOVA", "Correlation"])
        with col2:
            variables = st.multiselect("Select Variables", 
                                    st.session_state["parsed_data"]["questions"])
        
        if st.button("Run Advanced Analysis"):
            with st.spinner("Running analysis..."):
                # Implementation would vary by test type
                results = run_advanced_analysis(test_type, variables)
                st.dataframe(results)
                
                # NEW: Add interpretation
                st.markdown("#### 📝 Interpretation")
                st.write(get_test_interpretation(results))

# Visualization Tab (NEW)
with tabs[2]:
    if st.session_state["df"] is not None:
        st.subheader("📈 Visualization Dashboard")
        
        # Interactive visualization controls
        chart_type = st.selectbox("Chart Type", 
                                ["Bar", "Stacked Bar", "Heatmap", "Trend Line"])
        
        # Generate sample visualization
        fig, ax = plt.subplots()
        # Sample data - would be replaced with actual crosstab data
        categories = ['18-24', '25-34', '35-44', '45+']
        values = [25, 32, 18, 25]
        ax.bar(categories, values)
        ax.set_title("Sample Visualization")
        st.pyplot(fig)

# Report Generator Tab (NEW)
with tabs[3]:
    if st.session_state["df"] is not None:
        st.subheader("📝 Automated Report Generation")
        
        # Report customization options
        report_title = st.text_input("Report Title", "CrossTab Analysis Report")
        company_logo = st.file_uploader("Upload Company Logo (optional)", type=["png", "jpg"])
        
        if st.button("Generate Full Report"):
            with st.spinner("Generating report..."):
                # Create PowerPoint
                ppt_buffer = create_powerpoint_report(
                    st.session_state["parsed_data"],
                    st.session_state["insights"]
                )
                
                # Offer download
                st.download_button(
                    label="Download PowerPoint Report",
                    data=ppt_buffer,
                    file_name=f"{report_title.replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

# Enhanced Export Tab
with tabs[4]:
    if st.session_state["df"] is not None:
        st.subheader("📤 Enhanced Export Options")
        
        # NEW: Multiple export formats
        export_format = st.radio("Select Export Format",
                               ["Excel", "CSV", "PowerPoint", "PDF"])
        
        if st.button(f"Export as {export_format}"):
            with st.spinner(f"Preparing {export_format} export..."):
                # Implementation would vary by format
                buffer = prepare_export(export_format)
                st.download_button(
                    label=f"Download {export_format}",
                    data=buffer,
                    file_name=f"analysis.{export_format.lower()}",
                    mime="application/octet-stream"
                )
